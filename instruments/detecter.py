from pptx import Presentation

def detect_placeholders(location):
    """Detects the index of all the pages and placeholders in a ppt template.
    
    Loads a ppt template and add all the pages to it. Then fill each placeholder with its index and positions.
    Then saved as a ppt file as the result we get for the following development's need.
    """
    
    prs = Presentation(location)
    
    for layout in prs.slide_layouts:
      slide = prs.slides.add_slide(layout)
      
      for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        pos = (shape.left, shape.top, shape.width, shape.height)
        shape.text = str(idx) + '#' + str(pos)
    
    prs.save('DetectPlaceholders.pptx')


# Functions testing.

if __name__ == '__main__':
    
    detect_placeholders('../templates/Marxist.pptx')
    
    