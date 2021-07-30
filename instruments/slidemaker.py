from   pptx.dml.color   import RGBColor
from   pptx.util        import Pt

class BasicSlideMaker(object):
    """The basic slide maker class.
    
    All other slide makers are inherited from this class.
    
    Attributes:
        node: The OutlinNode object that this slide contains.
        prs_info: A dictory that contains the basic information of the presentation for content filling.
            title: A string contains the title of the presentation.
            date: A string contains the data of the presentation, usually formatted in "Date [day/month/year]".
            part_num: A string that indicates the serial number of the current part, usually formatted in "Part [num]".
            part_title: A string contains the title of the current part.
            color: A list that contains the main theme color in RGB format.
            author: A string contains the name of the author, usually formatted in "Report [name]".
        prs: The presentation object it works on.
        slide: The slide object it creates and works on.
    """
    
    def __init__(self, node, prs_info, prs, layout_idx):
        
        # Initialises the object's attributes.
        
        self.node = node
        self.prs_info = prs_info
        self.prs = prs
        
        # Creates a new slide, according to the layout index, in the presentation.
        
        layout = self.prs.slide_layouts[layout_idx]
        self.slide = self.prs.slides.add_slide(layout)


    def get_item(self, idx):
        """Returns the placeholder object according to the placeholder index"""
        return self.slide.shapes.placeholders[idx].text_frame.paragraphs[0]
    
    
class CoverSlideMaker(BasicSlideMaker):
    """The cover slide maker class.
    
    Noted that the slide index of cover slide is 0.
    
    Attributes:
        The attributes inherited from the BasicSlideMaker class aren't shown.
    """
    
    def __init__(self, node, prs_info, prs):
        
        BasicSlideMaker.__init__(self, node, prs_info, prs, 0)
        self.prs_info = prs_info

        self.title()
        self.date()
        self.abstract()

    def title(self):
        """Outputs the title.
        
        Noted that the index of title placeholder is 10.
        """
        
        item = self.get_item(10)
        item.text = self.prs_info['title']

    def date(self):
        """Outputs the date.
        
        Noted that the index of date placeholder is 19.
        """
        
        item = self.get_item(19)
        item.text = self.prs_info['date']

    def abstract(self):
        """Outputs the abstract.
        
        Noted that the index of abstract placeholder is 22.
        """
        
        item=self.get_item(22)
        num = 1
        
        for child in self.node.child:
            run = item.add_run()
            run.text = str(num) + '.' + child.title + '\n'
            num += 1

        run = item.add_run()
        run.text = '\n' + self.prs_info['author']
        
        
class SectionSlideMaker(BasicSlideMaker):
    """The section slide maker class.
    
    Noted that the index of first level section slide slide is 1, while the index of second level section slide is 2.
    
    Attributes:
        level: An integer that indicates the level of the section.
        The attributes inherited from the BasicSlideMaker class aren't shown.
    """
    
    def __init__(self, node, prs_info, prs, level):
        
        BasicSlideMaker.__init__(self, node, prs_info, prs, level)
        
        self.level = level
        self.header()
        self.section()

    def header(self):
        """Outputs the header.
        
        Noted that the index of header placeholder of first level section slice is 11.
        While the index of header placeholder of second level section slice is 13.
        """
        
        if self.level == 1:
            
            item = self.get_item(11)
            item.text = self.prs_info['title'] + '\n'

            run = item.add_run()
            run.text = 'Part ' + str(self.prs_info['part_num']) + ' '+self.prs_info['part_title']
            
        else:
            
            item = self.get_item(13)
            item.text = self.prs_info['title'] + '\n'

            run = item.add_run()
            run.text = 'Part ' + str(self.prs_info['part_num']) + ' '+self.prs_info['part_title']

    def section(self):
        """Outputs the section.
        
        Noted that the index of section placeholder of first level section slice is 13.
        While the index of section placeholder of second level section slice is 11.
        """
        
        if self.level == 1:
            
            item = self.get_item(13)
            item.text = '• ' + self.node.title + ' •'
      
            run = item.add_run()
            run.text = '••' + '\n'
            run.font.color.rgb = RGBColor(*self.prs_info['color'])
            
        else:
            
            item = self.get_item(11)
            item.text = '• ' + self.node.title + ' •'
      
            run = item.add_run()
            run.text = '••' + '\n'
            run.font.color.rgb = RGBColor(*[34,42,53]) # [34,42,53] is gray color the same as the second level section.
            
            
class TextSlideMaker(BasicSlideMaker):
    """The pure text slide maker class.
    
    Noted that the index of pure text slide slide is 3.
    
    Attributes:
        The attributes inherited from the BasicSlideMaker class aren't shown.
    """
    
    def __init__(self, node, prs_info, prs):
        
        BasicSlideMaker.__init__(self, node, prs_info, prs, 3)

        self.header()
        self.main_body()

    def header(self):
        """Outputs the header.
        
        Noted that the index of header placeholder is 11.
        """
        
        item = self.get_item(11)
        item.text = self.prs_info['title'] + '\n'

        run = item.add_run()
        run.text = 'Part ' + str(self.prs_info['part_num']) + ' '+self.prs_info['part_title']

    def main_body(self):
        """Outputs the main body.
        
        Noted that the index of header placeholder is 13.
        """
        
        item = self.get_item(13)
        
        item.text = '◤' + self.node.title + '\n'
        item.font.color.rgb = RGBColor(*self.prs_info['color'])
        
        run = item.add_run()
        run.text = self.node.note + '\n'
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(*[0,0,0]) # [0,0,0] is black.
        
        
class ImgTextSlideMaker(BasicSlideMaker):
    """The text with image slide maker class.
    
    Noted that the index of pure text slide slide is 4.
    
    Attributes:
        The attributes inherited from the BasicSlideMaker class aren't shown.
    """
    
    def __init__(self, node, prs_info, prs):
        
        BasicSlideMaker.__init__(self, node, prs_info, prs, 4)

        self.header()
        self.legend()
        self.main_body()

    def header(self):
        """Outputs the header.
        
        Noted that the index of header placeholder is 11.
        """
        
        item = self.get_item(11)
        item.text = self.prs_info['title'] + '\n'

        run = item.add_run()
        run.text = 'Part ' + str(self.prs_info['part_num']) + ' '+self.prs_info['part_title']

    def legend(self):
        """Outputs the image's legend.
        
        Noted that the index of header placeholder is 12.
        """
        
        item = self.get_item(12)
        item.text = '◤' + self.node.title

    def main_body(self):
        """Outputs the main body.
        
        Noted that the index of header placeholder is 13.
        """
        
        item = self.get_item(13)
        
        item.text = '◤' + self.node.title + '\n'
        item.font.color.rgb = RGBColor(*self.prs_info['color'])
        
        run = item.add_run()
        run.text = self.node.note + '\n'
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(*[0,0,0]) # [0,0,0] is black.
        
        
class BackCoverSlideMaker(BasicSlideMaker):
    """The back cover slide maker class.
    
    Noted that the index of back cover slide slide is 5.
    
    Attributes:
        The attributes inherited from the BasicSlideMaker class aren't shown.
    """
    
    def __init__(self, prs):
      
        BasicSlideMaker.__init__(self, None, None, prs, 5)
        
        